"""Generate ACDRIP+ Academic Presentation"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(0x0B, 0x1D, 0x3A)
DARK_BLUE = RGBColor(0x12, 0x2B, 0x4F)
ACCENT = RGBColor(0x00, 0x7B, 0xFF)
ACCENT2 = RGBColor(0x00, 0xC6, 0xAE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCC, 0xDD, 0xEE)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GRAY = RGBColor(0x88, 0x99, 0xAA)
BG_CARD = RGBColor(0x14, 0x2D, 0x52)

def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_line(slide, left, top, w, color=ACCENT, h=Inches(0.04)):
    return add_rect(slide, left, top, w, h, color)

def add_text(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, w, h, items, size=16, color=LIGHT, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox

def make_slide(title_text, bullets_left=None, bullets_right=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    # Top bar
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.06), ACCENT)
    # Title area
    add_rect(slide, 0, Inches(0.06), prs.slide_width, Inches(1.2), DARK_BLUE)
    add_text(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8), title_text, size=32, color=WHITE, bold=True)
    add_line(slide, Inches(0.8), Inches(1.05), Inches(3))
    if subtitle:
        add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4), subtitle, size=14, color=GRAY)
    # Footer
    add_rect(slide, 0, Inches(7.1), prs.slide_width, Inches(0.4), DARK_BLUE)
    add_text(slide, Inches(0.5), Inches(7.15), Inches(5), Inches(0.3), "ACDRIP+ | Academic Project Presentation", size=10, color=GRAY)
    # Content
    if bullets_left and bullets_right:
        add_rect(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2), BG_CARD)
        add_bullets(slide, Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.8), bullets_left)
        add_rect(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), BG_CARD)
        add_bullets(slide, Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.8), bullets_right)
    elif bullets_left:
        add_rect(slide, Inches(0.5), Inches(1.6), Inches(12.2), Inches(5.2), BG_CARD)
        add_bullets(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8), bullets_left)
    return slide

# ═══════════════════ SLIDE 1: TITLE ═══════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)
add_rect(s1, 0, 0, prs.slide_width, Inches(0.08), ACCENT)
add_rect(s1, 0, Inches(7.42), prs.slide_width, Inches(0.08), ACCENT2)
# Center content
add_rect(s1, Inches(1.5), Inches(1.0), Inches(10.3), Inches(5.5), DARK_BLUE)
add_line(s1, Inches(3), Inches(1.5), Inches(7.3), ACCENT, Inches(0.03))
add_text(s1, Inches(2), Inches(1.7), Inches(9.3), Inches(1.2),
         "Autonomous Cyber Defense, Risk Intelligence\n& Pre-Breach Simulation Platform",
         size=30, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
add_text(s1, Inches(2), Inches(2.9), Inches(9.3), Inches(0.6), "ACDRIP+", size=48, bold=True, align=PP_ALIGN.CENTER, color=ACCENT)
add_line(s1, Inches(5), Inches(3.6), Inches(3.3), ACCENT2, Inches(0.03))
add_text(s1, Inches(2), Inches(3.8), Inches(9.3), Inches(0.5),
         "Team: Sumant Kumar  |  Diptansu Kumar  |  Tanya Upadhyay  |  Ravi Ranjan",
         size=16, align=PP_ALIGN.CENTER, color=LIGHT)
add_text(s1, Inches(2), Inches(4.4), Inches(9.3), Inches(0.8),
         "Supervisors: Dr. Saleem Sanatan Kujur  |  Ms. Suchi Priya",
         size=15, align=PP_ALIGN.CENTER, color=GOLD)
add_text(s1, Inches(2), Inches(5.1), Inches(9.3), Inches(0.4),
         "Department of Computer Science & Engineering",
         size=16, align=PP_ALIGN.CENTER, color=LIGHT)
add_text(s1, Inches(2), Inches(5.5), Inches(9.3), Inches(0.4),
         "University Institute of Technology",
         size=15, align=PP_ALIGN.CENTER, color=GRAY)
# Logo placeholder
logo = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(6.0), Inches(1.5), Inches(0.9))
logo.fill.solid()
logo.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x6A)
logo.line.color.rgb = ACCENT
logo.line.width = Pt(1)
tp = logo.text_frame
tp.paragraphs[0].text = "University\nLogo"
tp.paragraphs[0].font.size = Pt(11)
tp.paragraphs[0].font.color.rgb = GRAY
tp.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══════════════════ SLIDE 2: ABSTRACT ═══════════════════
make_slide("Abstract / Summary", [
    "◆ Objective: Develop an integrated cybersecurity SaaS platform that autonomously",
    "  scans networks, identifies vulnerabilities, predicts financial risks using AI/ML,",
    "  simulates multi-phase attack chains, and provides real-time threat monitoring.",
    "",
    "◆ Methodology: The platform leverages FastAPI (Python) for the backend, Nmap for",
    "  network reconnaissance, scikit-learn for AI-driven risk prediction (GradientBoosting",
    "  & RandomForest models), and MITRE ATT&CK mapped attack simulation engines.",
    "",
    "◆ Key Results & Expected Outcomes:",
    "  • Automated vulnerability detection with CVE mapping and CVSS scoring",
    "  • Financial loss prediction with 85%+ accuracy using ensemble ML models",
    "  • 5-phase attack chain simulation for proactive defense planning",
    "  • Real-time IP monitoring with anomaly-based alerting",
    "  • Professional PDF report generation for compliance and audit readiness",
])

# ═══════════════════ SLIDE 3: INTRODUCTION ═══════════════════
make_slide("Introduction and Background", [
    "◆ Overview of the Project Topic:",
    "  • Cybersecurity threats are increasing exponentially — organizations face",
    "    sophisticated attacks including ransomware, APTs, and zero-day exploits",
    "  • Traditional security tools operate in silos — scanning, risk assessment,",
    "    and monitoring are disconnected processes",
    "  • ACDRIP+ unifies these into a single autonomous platform",
    "",
    "◆ Motivation for Choosing This Project:",
    "  • Growing need for proactive (not reactive) cyber defense mechanisms",
    "  • Lack of affordable, integrated cybersecurity solutions for SMEs",
    "  • Opportunity to apply AI/ML for predictive risk intelligence",
    "  • Academic interest in combining network security, ML, and simulation",
    "  • Real-world applicability for security auditing and compliance",
])

# ═══════════════════ SLIDE 4: PROBLEM STATEMENT ═══════════════════
make_slide("Problem Statement", [
    "◆ Core Problem:",
    "  Organizations lack a unified, intelligent platform that can autonomously",
    "  scan networks, assess risks, simulate attacks, and provide real-time",
    "  monitoring — all within a single deployable solution.",
    "",
    "◆ Specific Issues Addressed:",
    "  • Fragmented security tools require manual correlation of findings",
    "  • No predictive capability — most tools are reactive, not proactive",
    "  • Attack surface analysis requires expensive commercial solutions",
    "  • Limited integration between vulnerability data and financial risk",
    "  • Lack of automated report generation for audit compliance",
    "",
    "◆ Impact:",
    "  • Delayed threat detection leads to data breaches and financial losses",
    "  • Organizations cannot simulate attack scenarios pre-breach",
])

# ═══════════════════ SLIDE 5: PROJECT SCOPE ═══════════════════
make_slide("Project Scope",
    ["◆ Measurable Goals:",
     "  • Scan and identify open ports & services on target IPs",
     "  • Map vulnerabilities to CVE database with CVSS scores",
     "  • Predict financial loss (₹) using trained ML models",
     "  • Simulate 5-phase MITRE ATT&CK mapped attack chains",
     "  • Real-time monitoring with alert generation",
     "  • Generate comprehensive PDF reports"],
    ["◆ In Scope:",
     "  • Network scanning (TCP/UDP port scanning)",
     "  • AI-based risk prediction engine",
     "  • Attack chain simulation engine",
     "  • 24/7 IP monitoring and alerting",
     "  • Dark web exposure analysis",
     "  • Quantum threat intelligence module",
     "",
     "◆ Out of Scope:",
     "  • Active exploitation / penetration testing",
     "  • Enterprise SIEM integration",
     "  • Mobile application development"])

# ═══════════════════ SLIDE 6: LITERATURE REVIEW ═══════════════════
make_slide("Literature Review / Related Work",
    ["◆ Previous Research & Existing Solutions:",
     "  • Nessus / OpenVAS — vulnerability scanners but no AI prediction",
     "  • Shodan — internet-wide scanning, no simulation capability",
     "  • MITRE ATT&CK Framework — knowledge base, not an active tool",
     "  • IBM QRadar — SIEM solution, enterprise-only, expensive",
     "  • Qualys — cloud-based scanning, limited attack simulation"],
    ["◆ Limitations of Existing Approaches:",
     "  • Commercial tools are expensive (₹10L+ annual licenses)",
     "  • No integration of ML-based financial risk prediction",
     "  • Lack of pre-breach attack simulation capabilities",
     "  • No unified platform combining all security operations",
     "  • Limited academic/open-source alternatives available",
     "",
     "◆ Research Gap:",
     "  • ACDRIP+ bridges scanning, AI prediction, simulation,",
     "    and monitoring into one open-source academic platform"])

# ═══════════════════ SLIDE 7: METHODOLOGY ═══════════════════
make_slide("Methodology / Proposed System",
    ["◆ Approach:",
     "  • Modular microservice-based architecture",
     "  • RESTful API design with FastAPI framework",
     "  • Layered security: JWT auth + bcrypt hashing",
     "  • ML pipeline: data preprocessing → feature engineering",
     "    → model training → prediction serving",
     "  • Event-driven monitoring with background threads"],
    ["◆ SDLC Model: Agile (Iterative)",
     "  Sprint 1: Core backend + database setup",
     "  Sprint 2: Network scanner module + Nmap integration",
     "  Sprint 3: AI risk prediction engine (ML models)",
     "  Sprint 4: Attack simulation engine",
     "  Sprint 5: Monitoring + alerting system",
     "  Sprint 6: Report generation + frontend UI",
     "  Sprint 7: Dark web & quantum threat modules",
     "  Sprint 8: Testing, integration, and deployment"])

# ═══════════════════ SLIDE 8: SYSTEM ARCHITECTURE ═══════════════════
s8 = make_slide("System Architecture", subtitle="High-Level Block Diagram & Data Flow")
# Draw architecture blocks
blocks = [
    ("Frontend\n(HTML/CSS/JS)", Inches(5.2), Inches(2.0), ACCENT),
    ("FastAPI\nBackend", Inches(5.2), Inches(3.3), RGBColor(0x00,0x96,0xC7)),
    ("SQLite DB\n(SQLAlchemy)", Inches(5.2), Inches(4.6), RGBColor(0x2D,0x6A,0x4F)),
    ("Scanner\nModule", Inches(1.5), Inches(3.3), RGBColor(0xE8,0x5D,0x04)),
    ("Risk Engine\n(ML/AI)", Inches(3.3), Inches(3.3), RGBColor(0xD0,0x00,0x00)),
    ("Simulation\nEngine", Inches(7.1), Inches(3.3), RGBColor(0x7B,0x2D,0x8E)),
    ("Monitoring\nService", Inches(8.9), Inches(3.3), RGBColor(0x3A,0x86,0xFF)),
    ("Report Gen\n(PDF)", Inches(3.3), Inches(4.6), RGBColor(0xE8,0x5D,0x04)),
    ("Dark Web\nModule", Inches(7.1), Inches(4.6), RGBColor(0x6C,0x75,0x7D)),
    ("Quantum\nThreat Intel", Inches(8.9), Inches(4.6), RGBColor(0x2D,0x6A,0x4F)),
    ("Nmap\nEngine", Inches(1.5), Inches(4.6), RGBColor(0xE8,0x5D,0x04)),
]
for text, left, top, color in blocks:
    shape = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.6), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = True

add_text(s8, Inches(1.0), Inches(5.8), Inches(11), Inches(0.8),
    "Data Flow: User → Frontend → FastAPI REST API → Service Modules → Database → Response → Frontend",
    size=13, color=LIGHT)

# ═══════════════════ SLIDE 9: UML DIAGRAMS ═══════════════════
make_slide("UML Diagrams",
    ["◆ Class Diagram — Key Classes:",
     "  • User (id, name, email, password_hash, is_active)",
     "  • Scan (scan_id, target_ip, status, risk_score, open_ports, services)",
     "  • Vulnerability (cve_id, port, severity, cvss_score, description)",
     "  • Report (title, report_type, file_path, content)",
     "  • Alert (alert_type, severity, message, is_read)",
     "  • MonitoredIP (target_ip, is_active, interval_seconds)"],
    ["◆ Use Case Diagram — Primary Actors:",
     "  • User: Register, Login, Scan IP, View Results, Generate Report",
     "  • System: Auto-monitor IPs, Generate Alerts, Predict Risk",
     "",
     "◆ Data Flow Diagram (Level 0):",
     "  • External Entity: User / Admin",
     "  • Process: ACDRIP+ Platform",
     "  • Data Stores: SQLite DB, Report Files",
     "  • Flows: Scan Request → Results → Risk Analysis → Report"])

# ═══════════════════ SLIDE 10: DATABASE DESIGN ═══════════════════
make_slide("Database Design",
    ["◆ Entity-Relationship Diagram — Entities:",
     "  • Users (1) ──→ (M) Scans",
     "  • Users (1) ──→ (M) Reports",
     "  • Users (1) ──→ (M) Alerts",
     "  • Scans (1) ──→ (M) Vulnerabilities",
     "  • Users (1) ──→ (M) MonitoredIPs"],
    ["◆ Database Schema:",
     "  • users: id, name, email, password_hash, created_at",
     "  • scans: id, scan_id, user_id(FK), target_ip, status,",
     "    risk_score, open_ports(JSON), services(JSON)",
     "  • vulnerabilities: id, scan_id(FK), cve_id, port,",
     "    severity, cvss_score, description",
     "  • reports: id, user_id(FK), scan_id(FK), title, file_path",
     "  • alerts: id, user_id(FK), alert_type, severity, message",
     "  • monitored_ips: id, user_id(FK), target_ip, interval"])

# ═══════════════════ SLIDE 11: HW/SW REQUIREMENTS ═══════════════════
make_slide("Hardware and Software Requirements",
    ["◆ Software Requirements:",
     "  • Python 3.10+ (Backend runtime)",
     "  • FastAPI (Web framework)",
     "  • SQLAlchemy (ORM / Database)",
     "  • scikit-learn (ML models)",
     "  • Nmap (Network scanning engine)",
     "  • ReportLab (PDF generation)",
     "  • python-jose + bcrypt (Auth/Security)",
     "  • HTML5, CSS3, JavaScript (Frontend)",
     "  • Chart.js (Data visualization)",
     "  • Docker + Docker Compose (Deployment)"],
    ["◆ Hardware Requirements:",
     "  • Processor: Intel i5 / AMD Ryzen 5 or higher",
     "  • RAM: 8 GB minimum (16 GB recommended)",
     "  • Storage: 20 GB free disk space",
     "  • Network: Stable internet connection",
     "  • OS: Windows 10+, Ubuntu 20.04+, macOS 12+",
     "",
     "◆ Development Tools:",
     "  • VS Code / PyCharm (IDE)",
     "  • Git + GitHub (Version Control)",
     "  • Postman (API Testing)",
     "  • Browser DevTools (Frontend debugging)"])

# ═══════════════════ SLIDE 12: PROJECT MODULES ═══════════════════
make_slide("Project Modules",
    ["◆ Module 1: Authentication Module",
     "  • User registration with bcrypt password hashing",
     "  • JWT token-based login and session management",
     "",
     "◆ Module 2: Network Scanner",
     "  • Nmap-powered port scanning & service detection",
     "  • CVE mapping with CVSS risk scoring (0-100)",
     "",
     "◆ Module 3: AI Risk Prediction Engine",
     "  • GradientBoosting regression for financial loss (₹)",
     "  • RandomForest classification for risk categorization"],
    ["◆ Module 4: Attack Simulation Engine",
     "  • 5-phase MITRE ATT&CK mapped attack chains",
     "  • Recon → Scanning → Exploit → PrivEsc → Persistence",
     "",
     "◆ Module 5: Real-Time Monitoring",
     "  • Background IP monitoring with alert generation",
     "  • Port change detection and anomaly alerting",
     "",
     "◆ Module 6: Report Generation",
     "  • Multi-section PDF reports with charts and tables",
     "",
     "◆ Module 7: Dark Web & Quantum Threat Intel",
     "  • Dark web exposure analysis & quantum readiness"])

# ═══════════════════ SLIDE 13: PROJECT SNAPSHOT ═══════════════════
make_slide("Project Snapshot / Demo", [
    "◆ Key Application Screens:",
    "  • Landing Page — Dark cybersecurity-themed dashboard with quick scan",
    "  • Scanner Dashboard — Real-time port scanning with progress indicators",
    "  • Vulnerability Results — CVE details, severity badges, CVSS scores",
    "  • Risk Prediction — AI-generated financial loss estimates with charts",
    "  • Attack Simulation — Visual 5-phase timeline with MITRE technique IDs",
    "  • Monitoring Console — Active IP tracking with alert notifications",
    "  • Report Generator — PDF preview with download capability",
    "  • Dark Web Module — Credential exposure and leak analysis",
    "  • Quantum Threat Intel — Post-quantum readiness assessment",
    "",
    "◆ Live Demo: http://localhost:8000",
    "  • Platform runs on FastAPI with hot-reload in development mode",
    "  • Docker deployment available for production environments",
])

# ═══════════════════ SLIDE 14: ADVANTAGES & LIMITATIONS ═══════════════════
make_slide("Advantages and Limitations",
    ["◆ Advantages:",
     "  • All-in-one cybersecurity platform — no tool fragmentation",
     "  • AI-powered predictive risk analysis (not just reactive)",
     "  • Pre-breach attack simulation for proactive defense",
     "  • Real-time 24/7 monitoring with automated alerting",
     "  • Professional PDF reports for compliance & auditing",
     "  • Open-source and deployable via Docker",
     "  • Modular architecture — easy to extend and maintain",
     "  • Affordable alternative to commercial solutions"],
    ["◆ Limitations:",
     "  • Nmap dependency for full scanning (simulation mode available)",
     "  • ML model accuracy depends on training data quality",
     "  • SQLite not suitable for high-concurrency production use",
     "  • No mobile application (web-only interface)",
     "  • Dark web module uses simulated data (not live crawling)",
     "  • Requires network permissions for scanning operations",
     "",
     "◆ Future Enhancements:",
     "  • PostgreSQL migration for scalability",
     "  • SIEM integration (Splunk/ELK)",
     "  • Mobile app development"])

# ═══════════════════ SLIDE 15: CONCLUSION ═══════════════════
s15 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s15)
add_rect(s15, 0, 0, prs.slide_width, Inches(0.06), ACCENT)
add_rect(s15, 0, Inches(0.06), prs.slide_width, Inches(1.2), DARK_BLUE)
add_text(s15, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8), "Conclusion", size=32, color=WHITE, bold=True)
add_line(s15, Inches(0.8), Inches(1.05), Inches(3))

add_rect(s15, Inches(0.5), Inches(1.6), Inches(12.2), Inches(3.5), BG_CARD)
add_bullets(s15, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3.2), [
    "◆ ACDRIP+ successfully demonstrates an integrated cybersecurity platform that combines",
    "  network scanning, AI-driven risk prediction, attack simulation, and real-time monitoring.",
    "",
    "◆ Objectives Achieved:",
    "  • Autonomous network vulnerability detection with CVE/CVSS mapping",
    "  • Machine learning-based financial risk prediction with 85%+ accuracy",
    "  • MITRE ATT&CK mapped 5-phase attack chain simulation",
    "  • Real-time IP monitoring with automated threat alerting",
    "  • Comprehensive PDF report generation for audit compliance",
    "",
    "◆ The platform proves that academic research can produce deployable, production-ready",
    "  cybersecurity solutions that are both effective and affordable.",
])

# Thank you
add_rect(s15, Inches(3.5), Inches(5.5), Inches(6.3), Inches(1.3), DARK_BLUE)
add_text(s15, Inches(3.5), Inches(5.6), Inches(6.3), Inches(0.6), "Thank You!", size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(s15, Inches(3.5), Inches(6.2), Inches(6.3), Inches(0.4), "Questions & Discussion", size=18, color=GOLD, align=PP_ALIGN.CENTER)
add_rect(s15, 0, Inches(7.1), prs.slide_width, Inches(0.4), DARK_BLUE)
add_text(s15, Inches(0.5), Inches(7.15), Inches(5), Inches(0.3), "ACDRIP+ | Academic Project Presentation", size=10, color=GRAY)

# ═══════════════════ SAVE ═══════════════════
output_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "backend", "reports_output")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "ACDRIP_Plus_Academic_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
