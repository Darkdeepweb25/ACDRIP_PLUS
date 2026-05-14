from ppt_helpers import *
from pptx import Presentation
P=Presentation(); P.slide_width=SW; P.slide_height=SH

# S1: TITLE
s=P.slides.add_slide(P.slide_layouts[6]); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=NAVY
rect(s,0,0,SW,Inches(0.07),ACCENT); rect(s,0,Inches(7.43),SW,Inches(0.07),TEAL)
rect(s,Inches(1.8),Inches(1.2),Inches(9.7),Inches(5.0),DARK)
rect(s,Inches(3.5),Inches(1.6),Inches(6.3),Inches(0.03),ACCENT)
txt(s,Inches(2),Inches(1.8),Inches(9.3),Inches(1),"ACDRIP+",sz=52,c=ACCENT,bold=True,align=PP_ALIGN.CENTER)
txt(s,Inches(2),Inches(2.7),Inches(9.3),Inches(0.8),"Autonomous Cyber Defense, Risk Intelligence\n& Pre-Breach Simulation Platform",sz=22,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
rect(s,Inches(5.5),Inches(3.6),Inches(2.3),Inches(0.03),TEAL)
txt(s,Inches(2),Inches(3.8),Inches(9.3),Inches(0.4),"Sumant Kumar  •  Diptansu Kumar  •  Tanya Upadhyay  •  Ravi Ranjan",sz=14,c=LIGHT,align=PP_ALIGN.CENTER)
txt(s,Inches(2),Inches(4.3),Inches(9.3),Inches(0.4),"Supervisors:  Dr. Saleem Sanatan Kujur  |  Ms. Suchi Priya",sz=13,c=GOLD,align=PP_ALIGN.CENTER)
txt(s,Inches(2),Inches(4.9),Inches(9.3),Inches(0.3),"Department of Computer Science & Engineering",sz=14,c=LIGHT,align=PP_ALIGN.CENTER)
txt(s,Inches(2),Inches(5.3),Inches(9.3),Inches(0.3),"University Institute of Technology",sz=13,c=GRAY,align=PP_ALIGN.CENTER)
sh=rrect(s,Inches(6.1),Inches(5.7),Inches(1.2),Inches(0.7),RGBColor(0x1A,0x3A,0x6A),text="Logo",tsz=11,tc=GRAY,border=ACCENT)

# S2: ABSTRACT
s=slide_base(P,"Abstract / Summary",2)
rect(s,Inches(0.5),Inches(1.5),Inches(12.2),Inches(5.3),CARD)
bullets(s,Inches(0.9),Inches(1.7),Inches(11.4),Inches(4.8),[
 "▸ Objective: Build a unified cybersecurity platform for autonomous scanning, AI risk prediction, and attack simulation",
 "▸ Methodology: FastAPI backend + Nmap scanning + scikit-learn ML + MITRE ATT&CK simulation",
 "▸ Outcome: Automated vulnerability detection, financial loss prediction, 5-phase attack chain simulation, and PDF reporting"],sz=16)

# S3: INTRODUCTION
s=slide_base(P,"Introduction & Background",3)
rect(s,Inches(0.5),Inches(1.5),Inches(5.8),Inches(5.3),CARD)
txt(s,Inches(0.9),Inches(1.6),Inches(5),Inches(0.4),"Overview",sz=18,c=ACCENT,bold=True)
bullets(s,Inches(0.9),Inches(2.1),Inches(5),Inches(2.5),[
 "▸ Cyber threats growing exponentially",
 "▸ Traditional tools operate in silos",
 "▸ ACDRIP+ unifies scanning, risk analysis, simulation & monitoring"],sz=14)
rect(s,Inches(6.8),Inches(1.5),Inches(5.8),Inches(5.3),CARD)
txt(s,Inches(7.2),Inches(1.6),Inches(5),Inches(0.4),"Motivation",sz=18,c=TEAL,bold=True)
bullets(s,Inches(7.2),Inches(2.1),Inches(5),Inches(2.5),[
 "▸ Need for proactive cyber defense",
 "▸ Lack of affordable integrated solutions",
 "▸ Apply AI/ML for predictive intelligence",
 "▸ Real-world security auditing use case"],sz=14)

# S4: PROBLEM STATEMENT
s=slide_base(P,"Problem Statement",4)
rect(s,Inches(0.5),Inches(1.5),Inches(12.2),Inches(5.3),CARD)
txt(s,Inches(0.9),Inches(1.7),Inches(11),Inches(0.4),"Core Problem",sz=20,c=RED,bold=True)
bullets(s,Inches(0.9),Inches(2.3),Inches(11),Inches(1.5),[
 "▸ No unified platform that autonomously scans, assesses risks, simulates attacks & monitors in real-time"],sz=16)
txt(s,Inches(0.9),Inches(3.2),Inches(11),Inches(0.4),"Key Issues",sz=18,c=ORANGE,bold=True)
bullets(s,Inches(0.9),Inches(3.7),Inches(11),Inches(2.5),[
 "▸ Fragmented tools require manual correlation",
 "▸ Reactive approach — no predictive capability",
 "▸ Commercial solutions cost ₹10L+ annually",
 "▸ No integration between vulnerability data and financial risk"],sz=14)

# S5: PROJECT SCOPE
s=slide_base(P,"Project Scope",5)
rect(s,Inches(0.5),Inches(1.5),Inches(5.8),Inches(5.3),CARD)
txt(s,Inches(0.9),Inches(1.6),Inches(5),Inches(0.4),"✓  In Scope",sz=18,c=GREEN,bold=True)
bullets(s,Inches(0.9),Inches(2.1),Inches(5),Inches(4),[
 "▸ Network port scanning & service detection",
 "▸ AI-based financial risk prediction",
 "▸ MITRE ATT&CK attack simulation",
 "▸ 24/7 IP monitoring & alerting",
 "▸ Dark web exposure analysis",
 "▸ Quantum threat intelligence",
 "▸ PDF report generation"],sz=14)
rect(s,Inches(6.8),Inches(1.5),Inches(5.8),Inches(5.3),CARD)
txt(s,Inches(7.2),Inches(1.6),Inches(5),Inches(0.4),"✗  Out of Scope",sz=18,c=RED,bold=True)
bullets(s,Inches(7.2),Inches(2.1),Inches(5),Inches(4),[
 "▸ Active exploitation / pen testing",
 "▸ Enterprise SIEM integration",
 "▸ Mobile application",
 "▸ Live dark web crawling"],sz=14)

# S6: LITERATURE REVIEW
s=slide_base(P,"Literature Review",6)
rect(s,Inches(0.5),Inches(1.5),Inches(12.2),Inches(5.3),CARD)
txt(s,Inches(0.9),Inches(1.6),Inches(5),Inches(0.4),"Existing Solutions",sz=18,c=ACCENT,bold=True)
# Table-like layout
for i,(tool,desc,lim) in enumerate([("Nessus/OpenVAS","Vulnerability Scanner","No AI prediction"),
 ("Shodan","Internet Scanner","No simulation"),("IBM QRadar","SIEM","Expensive, enterprise-only"),
 ("Qualys","Cloud Scanner","Limited simulation")]):
    y=Inches(2.2+i*0.7)
    rrect(s,Inches(0.9),y,Inches(2),Inches(0.5),DARK,text=tool,tsz=11,tc=ACCENT,border=ACCENT)
    txt(s,Inches(3.1),y+Inches(0.05),Inches(3.5),Inches(0.4),desc,sz=12,c=LIGHT)
    txt(s,Inches(6.8),y+Inches(0.05),Inches(5),Inches(0.4),"⚠ "+lim,sz=12,c=ORANGE)
txt(s,Inches(0.9),Inches(5.2),Inches(11),Inches(0.4),"Research Gap: No open-source platform combines scanning + AI prediction + simulation + monitoring",sz=14,c=GOLD,bold=True)

# S7: METHODOLOGY
s=slide_base(P,"Methodology / Proposed System",7)
rect(s,Inches(0.5),Inches(1.5),Inches(5.8),Inches(5.3),CARD)
txt(s,Inches(0.9),Inches(1.6),Inches(5),Inches(0.4),"Approach",sz=18,c=ACCENT,bold=True)
bullets(s,Inches(0.9),Inches(2.1),Inches(5),Inches(3),[
 "▸ Modular microservice architecture",
 "▸ RESTful API with FastAPI",
 "▸ JWT auth + bcrypt security",
 "▸ ML pipeline with scikit-learn"],sz=14)
# SDLC diagram
rect(s,Inches(6.8),Inches(1.5),Inches(5.8),Inches(5.3),CARD)
txt(s,Inches(7.2),Inches(1.6),Inches(5),Inches(0.4),"Agile SDLC Sprints",sz=18,c=TEAL,bold=True)
sprints=[("S1","Core Setup",ACCENT),("S2","Scanner",ORANGE),("S3","AI/ML",RED),("S4","Simulation",PURPLE),
 ("S5","Monitoring",GREEN),("S6","Reports+UI",TEAL),("S7","Dark/Quantum",PINK),("S8","Testing",GOLD)]
for i,(sp,desc,clr) in enumerate(sprints):
    r,c2=i//4,i%4
    x,y=Inches(7.3+c2*1.3),Inches(2.2+r*1.5)
    rrect(s,x,y,Inches(1.1),Inches(0.6),clr,text=sp,tsz=12)
    txt(s,x-Inches(0.1),y+Inches(0.65),Inches(1.3),Inches(0.4),desc,sz=9,c=LIGHT,align=PP_ALIGN.CENTER)

# S8: SYSTEM ARCHITECTURE (DIAGRAM)
s=slide_base(P,"System Architecture",8)
# Frontend layer
rrect(s,Inches(4.5),Inches(1.5),Inches(4.3),Inches(0.8),ACCENT,text="Frontend  (HTML / CSS / JS / Chart.js)",tsz=13)
arrow_d(s,Inches(6.5),Inches(2.3))
# API layer
rrect(s,Inches(3.5),Inches(2.8),Inches(6.3),Inches(0.8),RGBColor(0x00,0x96,0xC7),text="FastAPI Backend  (REST API + JWT Auth)",tsz=13)
# Module boxes
mods=[("Scanner\n(Nmap)",Inches(0.5),ORANGE),("Risk Engine\n(ML/AI)",Inches(2.7),RED),
 ("Simulation\n(ATT&CK)",Inches(4.9),PURPLE),("Monitoring\n(24/7)",Inches(7.1),GREEN),
 ("Reports\n(PDF)",Inches(9.3),TEAL),("DarkWeb &\nQuantum",Inches(11.3),PINK)]
arrow_d(s,Inches(6.5),Inches(3.6))
for label,x,clr in mods:
    rrect(s,x,Inches(4.3),Inches(1.8),Inches(0.9),clr,text=label,tsz=11)
# DB layer
arrow_d(s,Inches(6.5),Inches(5.2))
rrect(s,Inches(4.0),Inches(5.7),Inches(5.3),Inches(0.8),RGBColor(0x2D,0x6A,0x4F),text="SQLite Database  (SQLAlchemy ORM)",tsz=13)
txt(s,Inches(0.5),Inches(6.7),Inches(12),Inches(0.3),"Data Flow:  User → Frontend → REST API → Service Module → Database → Response",sz=12,c=GRAY)

P.save("__ppt_temp.pptx")
print("Part1 done: 8 slides")
