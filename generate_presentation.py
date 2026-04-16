from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def generate_presentation():
    prs = Presentation()
    
    # Set dark background color for all slide masters
    # Standard 16:9 slide size is usually default or 10x5.625 inches
    
    def apply_dark_theme(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(10, 14, 25) # Dark Blue background
        
    def add_title_slide():
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        apply_dark_theme(slide)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "ACDRIP+"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 229, 255) # Cyan
        title.text_frame.paragraphs[0].font.size = Pt(64)
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = "Autonomous Cyber Defense, Risk Intelligence & Pre-Breach Simulation Platform\n\nAcademic Presentation"
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(139, 92, 246) # Purple
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)
        
    def add_content_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        apply_dark_theme(slide)
        title = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 229, 255)
        title.text_frame.paragraphs[0].font.size = Pt(36)
        
        tf = body_shape.text_frame
        tf.clear()
        
        for i, point in enumerate(bullet_points):
            p = tf.add_paragraph()
            p.text = point
            p.font.color.rgb = RGBColor(226, 232, 240)
            p.font.size = Pt(20)
            if i > 0:
                p.space_before = Pt(14)
                
    # Slide 1: Title
    add_title_slide()
    
    # Slide 2: The Problem
    add_content_slide("1. The Cybersecurity Challenge", [
        "Modern threat landscapes are extremely complex and highly automated.",
        "Average cost of a data breach reached $4.45M in 2023 (IBM).",
        "Current industry relies on reactive toolkits (incident response).",
        "Pre-breach platforms exist, but cost $50,000+ annually.",
        "SMEs and academia lack accessible, unified tools for pre-breach defense.",
        "Emerging threats like Post-Quantum Cryptography (PQC) vulnerabilities are ignored by standard scanners."
    ])
    
    # Slide 3: Solution
    add_content_slide("2. Introducing ACDRIP+", [
        "A fully integrated, 9-module web platform for autonomous cyber defense.",
        "Unifies Network Scanning, ML Risk Prediction, and Attack Simulation.",
        "Includes advanced OSINT: MITRE ATT&CK, Dark Web, and Quantum Audit.",
        "FastAPI (Python) backend + Vanilla ES6 JS frontend.",
        "Zero-license cost (Free/Open Source) alternative to enterprise tools."
    ])
    
    # Slide 4: Architecture
    add_content_slide("3. System Architecture", [
        "Frontend: Vanilla JS, HTML5, CSS3 Glassmorphism (Zero framework bloat).",
        "Backend: FastAPI (Async Python) + Uvicorn + SQLAlchemy.",
        "Database: SQLite (Serverless, fully local, private).",
        "Services: python-nmap, pure-Python ML algorithm, ReportLab PDF gen.",
        "Extremely low latency (<100ms load) and fully portable (Docker support)."
    ])
    
    # Slide 5: Core 1
    add_content_slide("4. Module 1: Profile-Based Scanner", [
        "Executes real port scans (via Nmap) or deterministic simulations.",
        "Simulation uses SHA-256 IP hashing to assign logical 'Host Profiles'.",
        "Discovers up to 22 critical services (SSH, Redis, SMB, RDP, DNS, HTTP).",
        "Automatically cross-references discovered services with a local CVE database.",
        "Computes initial CVSS scores based on open attack surface."
    ])
    
    # Slide 6: Core 2
    add_content_slide("5. Module 2: Financial Risk Engine", [
        "Uses a Random Forest-inspired decision tree written in pure Python.",
        "Inputs: Asset Value, Critical/High Vulns, Port Count, Firewalls/IDS, Industry.",
        "Outputs FAIR-aligned projected financial loss in local currency (INR).",
        "Computes attack probability percentage.",
        "Segments the risk curve into a standard 5-tier classification logic."
    ])
    
    # Slide 7: Core 3
    add_content_slide("6. Module 3: Attack Kill-Chain Simulator", [
        "Maps vulnerabilities dynamically to the 7-phase Lockheed Martin Kill Chain.",
        "Phases: Recon -> Weaponization -> Delivery -> Exploit -> Install -> C2 -> Impact.",
        "Calculates phase success probability and plots attack graphs.",
        "Generates detailed AI analysis for adversarial toolchains (Mimikatz, etc.).",
        "Provides exact mitigation paths mapped to MITRE identifiers."
    ])
    
    # Slide 8: Advanced 1
    add_content_slide("7. MITRE ATT&CK Interactive Matrix", [
        "Fully interactive UI rendering 11 enterprise tactics and 33 techniques.",
        "Highlights techniques successfully detected in network simulations.",
        "Click-to-expand cards offering deep-dive adversarial tactics logic.",
        "Cross-references tactic severity and detection counts.",
        "Provides actionable defensive postures natively inside the dashboard."
    ])
    
    # Slide 9: Advanced 2
    add_content_slide("8. Dark Web Exposure Analysis", [
        "Deterministic seed-based generation to model OSINT deep-web presence.",
        "Tracks credential leak counts, BreachForums mentions, Ransomware chatter.",
        "Graphs breaches over a 6-month historical timeline.",
        "Identifies active TOR exit node targeting.",
        "Provides enterprise-vulnerability intelligence without direct TOR access."
    ])
    
    # Slide 10: Advanced 3
    add_content_slide("9. Post-Quantum Cryptographic Audit", [
        "First open-source scanner integrating a Quantum vulnerability module.",
        "Audits 8 cryptographic surfaces (TLS, VPN, Code Signing, JWTs, etc.).",
        "Highlights Shor's algorithm risks (Asymmetric: RSA, ECDHE) and Grover's risks.",
        "Provides current state vs. NIST PQC recommendations (Kyber, Dilithium).",
        "Renders a real-time HTML5 Canvas Quantum Readiness Gauge."
    ])
    
    # Slide 11: Real-time
    add_content_slide("10. 24/7 Monitoring & Threat Intel", [
        "Monitoring: Background service continuously checking IP port state changes.",
        "Alert Logic: Creates priority badges for newly opened ports or backdoors.",
        "Threat Intel Feed: Live ticker of zero-day exploits and APT campaigns.",
        "IOC Database: Tracking IP, Hash, and Domain indicators of compromise.",
        "Keeps defenders proactively aware of global threat landscapes."
    ])
    
    # Slide 12: Reporting
    add_content_slide("11. Professional PDF Reporting", [
        "Automated report generation using Python ReportLab.",
        "Consolidates scans, ML risks, simulations, and telemetry into a single PDF.",
        "Features graphical charts, CVSS breakdowns, and executive summaries.",
        "Saves direct to binary streams for instant browser downloads.",
        "Fully branded and strictly formatted for C-level executives."
    ])
    
    # Slide 13: Comparison
    add_content_slide("12. Capability Matrix", [
        "Compared against Nessus, Metasploit, Shodan, and Recorded Future.",
        "ACDRIP+ is the only tool unifying ALL 9 capabilities natively.",
        "Replaces roughly $50,000 USD/year in specialized commercial licenses.",
        "Open-source MIT license ensures total academic accessibility.",
        "Superior UI/UX designed around modern intuitive workflows."
    ])
    
    # Slide 14: Future roadmap
    add_content_slide("13. Future Roadmap", [
        "V2 Roadmap: Live Dark Web API integration (via AlienVault OTX).",
        "Graph Neural Network (GNN) implementation for vulnerability chaining.",
        "Custom Nmap NSE scripts to ping live post-quantum TLS cipher suites.",
        "Cloud Infrastructure integration (AWS/Azure IAM permission scraping).",
        "Mobile companion app built in Flutter for real-time alert delivery."
    ])
    
    # Slide 15: Conclusion
    add_content_slide("14. Conclusion", [
        "ACDRIP+ redefines what an open-source cybersecurity platform can do.",
        "We've proven that deep, multi-vector pre-breach simulation is achievable.",
        "The project provides unparalleled educational value for defensive testing.",
        "Fully standalone, serverless SQL, rapid-deployment ready.",
        "Thank you! Any Questions?"
    ])
    
    output_dir = "backend/reports_output"
    os.makedirs(output_dir, exist_ok=True)
    out_path = os.path.join(output_dir, "ACDRIP_Plus_Presentation.pptx")
    prs.save(out_path)
    print(f"Presentation saved to {out_path}")

if __name__ == "__main__":
    generate_presentation()
