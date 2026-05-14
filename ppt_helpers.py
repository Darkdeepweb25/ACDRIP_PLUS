"""PPT Helper Functions & Design System"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Design Tokens ──
NAVY = RGBColor(0x0A, 0x1A, 0x35)
DARK = RGBColor(0x0F, 0x23, 0x44)
CARD = RGBColor(0x12, 0x2C, 0x55)
ACCENT = RGBColor(0x00, 0x7B, 0xFF)
TEAL = RGBColor(0x00, 0xC6, 0xAE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xBB, 0xCC, 0xDD)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GRAY = RGBColor(0x77, 0x88, 0x99)
RED = RGBColor(0xFF, 0x45, 0x45)
GREEN = RGBColor(0x00, 0xC8, 0x53)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)
PINK = RGBColor(0xE9, 0x1E, 0x63)
SW = Inches(13.333)
SH = Inches(7.5)

def rect(slide, l, t, w, h, color, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    if border: s.line.color.rgb = border; s.line.width = Pt(1)
    return s

def rrect(slide, l, t, w, h, color, text="", tsz=10, tc=WHITE, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    if border: s.line.color.rgb = border; s.line.width = Pt(1.5)
    if text:
        tf = s.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(tsz); p.font.color.rgb = tc
        p.font.bold = True; p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    return s

def txt(slide, l, t, w, h, text, sz=18, c=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = bold; p.font.name = "Calibri"; p.alignment = align
    return tb

def bullets(slide, l, t, w, h, items, sz=15, c=LIGHT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = "Calibri"; p.space_after = Pt(8)
    return tb

def arrow_r(slide, l, t, w, h=Inches(0.3), color=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def arrow_d(slide, l, t, w=Inches(0.25), h=Inches(0.4), color=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def circle(slide, l, t, sz, color, text="", tsz=9, tc=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    if text:
        tf = s.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(tsz); p.font.color.rgb = tc
        p.font.bold = True; p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER
    return s

def diamond(slide, l, t, w, h, color, text="", tsz=9, tc=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    if text:
        tf = s.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(tsz); p.font.color.rgb = tc
        p.font.bold = True; p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER
    return s

def slide_base(prs, title, slide_num=None):
    """Create a slide with standard header/footer"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    rect(sl, 0, 0, SW, Inches(0.05), ACCENT)
    rect(sl, 0, Inches(0.05), SW, Inches(1.1), DARK)
    txt(sl, Inches(0.7), Inches(0.15), Inches(10), Inches(0.7), title, sz=30, c=WHITE, bold=True)
    rect(sl, Inches(0.7), Inches(0.9), Inches(2.5), Inches(0.04), ACCENT)
    rect(sl, 0, Inches(7.15), SW, Inches(0.35), DARK)
    footer = f"ACDRIP+  •  Slide {slide_num}" if slide_num else "ACDRIP+"
    txt(sl, Inches(0.5), Inches(7.18), Inches(4), Inches(0.25), footer, sz=9, c=GRAY)
    return sl

def icon_box(slide, l, t, icon_text, label, color):
    """Small icon-style box with label below"""
    rrect(slide, l, t, Inches(0.7), Inches(0.7), color, text=icon_text, tsz=20, tc=WHITE)
    txt(slide, l - Inches(0.15), t + Inches(0.75), Inches(1.0), Inches(0.3), label, sz=9, c=LIGHT, align=PP_ALIGN.CENTER)
