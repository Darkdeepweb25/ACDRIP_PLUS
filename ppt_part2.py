from pptx import Presentation
from ppt_helpers import *
P=Presentation("__ppt_temp.pptx")

# S9: UML - CLASS DIAGRAM
s=slide_base(P,"UML — Class Diagram",9)
classes=[("User",["id: str","name: str","email: str","password_hash"],Inches(0.5),Inches(1.7),ACCENT),
 ("Scan",["scan_id: str","target_ip: str","risk_score: float","open_ports: JSON"],Inches(4.0),Inches(1.7),ORANGE),
 ("Vulnerability",["cve_id: str","severity: str","cvss_score: float"],Inches(7.8),Inches(1.7),RED),
 ("Report",["title: str","report_type: str","file_path: str"],Inches(0.5),Inches(4.5),TEAL),
 ("Alert",["alert_type: str","severity: str","message: str"],Inches(4.0),Inches(4.5),PURPLE),
 ("MonitoredIP",["target_ip: str","is_active: bool","interval: int"],Inches(7.8),Inches(4.5),GREEN)]
for name,attrs,x,y,clr in classes:
    rect(s,x,y,Inches(3.2),Inches(0.45),clr)
    txt(s,x+Inches(0.1),y+Inches(0.05),Inches(3),Inches(0.35),name,sz=13,c=WHITE,bold=True)
    rect(s,x,y+Inches(0.45),Inches(3.2),Inches(1.6),CARD)
    tb=s.shapes.add_textbox(x+Inches(0.15),y+Inches(0.5),Inches(2.9),Inches(1.5))
    tf=tb.text_frame; tf.word_wrap=True
    for i,a in enumerate(attrs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=a; p.font.size=Pt(10); p.font.color.rgb=LIGHT; p.font.name="Consolas"
# Relationship arrows
arrow_r(s,Inches(3.7),Inches(2.1),Inches(0.3))
arrow_r(s,Inches(7.2),Inches(2.1),Inches(0.6))
arrow_d(s,Inches(1.8),Inches(3.8),h=Inches(0.5))
arrow_d(s,Inches(5.3),Inches(3.8),h=Inches(0.5))

# S10: USE CASE DIAGRAM
s=slide_base(P,"UML — Use Case & DFD",10)
# Use case side
rect(s,Inches(0.4),Inches(1.4),Inches(6),Inches(5.4),CARD)
txt(s,Inches(0.7),Inches(1.5),Inches(3),Inches(0.35),"Use Case Diagram",sz=16,c=ACCENT,bold=True)
# Actor
rrect(s,Inches(0.7),Inches(2.5),Inches(1.2),Inches(0.6),DARK,text="👤 User",tsz=11,tc=LIGHT,border=ACCENT)
cases=["Register / Login","Scan Network","View Results","Predict Risk","Simulate Attack","Generate Report"]
for i,uc in enumerate(cases):
    y=Inches(2.1+i*0.55)
    rrect(s,Inches(2.8),y,Inches(3.2),Inches(0.45),DARK,text=uc,tsz=10,tc=LIGHT,border=TEAL)
# DFD side
rect(s,Inches(6.8),Inches(1.4),Inches(6),Inches(5.4),CARD)
txt(s,Inches(7.1),Inches(1.5),Inches(3),Inches(0.35),"Data Flow Diagram",sz=16,c=TEAL,bold=True)
rrect(s,Inches(7.2),Inches(2.3),Inches(1.6),Inches(0.7),DARK,text="User",tsz=11,tc=LIGHT,border=ACCENT)
arrow_r(s,Inches(8.8),Inches(2.5),Inches(0.5))
circle(s,Inches(9.4),Inches(2.1),Inches(1.1),ACCENT,text="ACDRIP+\nPlatform",tsz=9)
arrow_r(s,Inches(10.5),Inches(2.5),Inches(0.5))
rrect(s,Inches(11.1),Inches(2.3),Inches(1.4),Inches(0.7),DARK,text="Reports",tsz=10,tc=LIGHT,border=TEAL)
arrow_d(s,Inches(9.85),Inches(3.2),h=Inches(0.4))
rrect(s,Inches(9.1),Inches(3.7),Inches(1.8),Inches(0.6),RGBColor(0x2D,0x6A,0x4F),text="SQLite DB",tsz=10,border=GREEN)

# S11: DATABASE DESIGN (ER DIAGRAM)
s=slide_base(P,"Database Design — ER Diagram",11)
entities=[("Users",Inches(5.5),Inches(1.6),ACCENT),("Scans",Inches(1.5),Inches(3.5),ORANGE),
 ("Vulns",Inches(1.5),Inches(5.5),RED),("Reports",Inches(5.5),Inches(5.5),TEAL),
 ("Alerts",Inches(9.5),Inches(3.5),PURPLE),("MonitoredIPs",Inches(9.5),Inches(5.5),GREEN)]
for name,x,y,clr in entities:
    rrect(s,x,y,Inches(2.2),Inches(0.8),clr,text=name,tsz=14)
# Relationships
rels=[("1:M",Inches(3.5),Inches(2.7)),("1:M",Inches(7.5),Inches(2.7)),("1:M",Inches(5.5),Inches(4.5)),
 ("1:M",Inches(1.5),Inches(4.8))]
for label,x,y in rels:
    diamond(s,x,y,Inches(0.9),Inches(0.6),GOLD,text=label,tsz=9,tc=NAVY)
txt(s,Inches(0.5),Inches(6.6),Inches(12),Inches(0.3),"Primary keys: UUID  •  Foreign keys link Users→Scans→Vulnerabilities  •  JSON fields for ports, services",sz=11,c=GRAY)

# S12: HW/SW REQUIREMENTS
s=slide_base(P,"Hardware & Software Requirements",12)
rect(s,Inches(0.5),Inches(1.5),Inches(5.8),Inches(5.3),CARD)
txt(s,Inches(0.9),Inches(1.6),Inches(5),Inches(0.35),"Software Stack",sz=18,c=ACCENT,bold=True)
sw=[("FastAPI","Backend Framework"),("SQLAlchemy","ORM / Database"),("scikit-learn","ML Models"),
 ("Nmap","Network Scanner"),("ReportLab","PDF Generation"),("Docker","Containerization")]
for i,(t,d) in enumerate(sw):
    y=Inches(2.15+i*0.55)
    rrect(s,Inches(0.9),y,Inches(1.8),Inches(0.42),DARK,text=t,tsz=10,tc=ACCENT,border=ACCENT)
    txt(s,Inches(2.9),y+Inches(0.05),Inches(3),Inches(0.35),d,sz=12,c=LIGHT)
rect(s,Inches(6.8),Inches(1.5),Inches(5.8),Inches(5.3),CARD)
txt(s,Inches(7.2),Inches(1.6),Inches(5),Inches(0.35),"Hardware",sz=18,c=TEAL,bold=True)
bullets(s,Inches(7.2),Inches(2.2),Inches(5),Inches(3),[
 "▸ CPU: Intel i5 / Ryzen 5+","▸ RAM: 8 GB minimum","▸ Storage: 20 GB free",
 "▸ Network: Stable connection","▸ OS: Win 10+ / Ubuntu 20.04+"],sz=14)

# S13: PROJECT MODULES
s=slide_base(P,"Project Modules",13)
mods2=[("🔐","Auth",ACCENT),("🔍","Scanner",ORANGE),("🧠","AI Risk",RED),
 ("⚔","Simulation",PURPLE),("📡","Monitor",GREEN),("📄","Reports",TEAL),("🌐","DarkWeb\nQuantum",PINK)]
for i,(ico,name,clr) in enumerate(mods2):
    x=Inches(0.5+i*1.8)
    rrect(s,x,Inches(1.7),Inches(1.5),Inches(1.2),clr,text=ico,tsz=28)
    txt(s,x,Inches(2.95),Inches(1.5),Inches(0.5),name,sz=11,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
bullets(s,Inches(0.7),Inches(3.7),Inches(12),Inches(3),[
 "▸ Auth: JWT tokens + bcrypt hashing  |  Scanner: Nmap port scanning + CVE mapping",
 "▸ Risk Engine: GradientBoosting + RandomForest for financial loss prediction (₹)",
 "▸ Simulation: 5-phase MITRE ATT&CK chain (Recon → Exploit → Persistence)",
 "▸ Monitoring: Background thread IP monitoring + anomaly alerting",
 "▸ Reports: Multi-section PDF with charts  |  DarkWeb & Quantum threat modules"],sz=13)

# S14: ADVANTAGES & LIMITATIONS
s=slide_base(P,"Advantages & Limitations",14)
rect(s,Inches(0.5),Inches(1.5),Inches(5.8),Inches(5.3),CARD)
txt(s,Inches(0.9),Inches(1.6),Inches(5),Inches(0.35),"✓  Advantages",sz=18,c=GREEN,bold=True)
bullets(s,Inches(0.9),Inches(2.1),Inches(5),Inches(4),[
 "▸ Unified all-in-one platform","▸ AI-powered predictive analysis",
 "▸ Pre-breach attack simulation","▸ Real-time 24/7 monitoring",
 "▸ Open-source & Docker-ready","▸ Modular & extensible"],sz=14)
rect(s,Inches(6.8),Inches(1.5),Inches(5.8),Inches(5.3),CARD)
txt(s,Inches(7.2),Inches(1.6),Inches(5),Inches(0.35),"✗  Limitations",sz=18,c=RED,bold=True)
bullets(s,Inches(7.2),Inches(2.1),Inches(5),Inches(4),[
 "▸ Nmap dependency for full scanning","▸ SQLite limits concurrency",
 "▸ No mobile app","▸ Simulated dark web data",
 "▸ ML accuracy depends on data quality"],sz=14)

# S15: CONCLUSION
s=slide_base(P,"Conclusion",15)
rect(s,Inches(0.5),Inches(1.5),Inches(12.2),Inches(3.2),CARD)
bullets(s,Inches(0.9),Inches(1.7),Inches(11),Inches(2.8),[
 "▸ ACDRIP+ successfully delivers a unified cybersecurity platform combining scanning, AI, simulation & monitoring",
 "▸ All primary objectives achieved: vulnerability detection, risk prediction, attack simulation, real-time alerting",
 "▸ The platform demonstrates that academic research can produce deployable, production-ready solutions"],sz=15)
rect(s,Inches(3.5),Inches(5.0),Inches(6.3),Inches(1.5),DARK)
txt(s,Inches(3.5),Inches(5.1),Inches(6.3),Inches(0.7),"Thank You!",sz=40,c=ACCENT,bold=True,align=PP_ALIGN.CENTER)
txt(s,Inches(3.5),Inches(5.8),Inches(6.3),Inches(0.5),"Questions & Discussion",sz=18,c=GOLD,align=PP_ALIGN.CENTER)

import os
out=os.path.join(os.path.dirname(os.path.abspath(__file__)),"backend","reports_output")
os.makedirs(out,exist_ok=True)
fp=os.path.join(out,"ACDRIP_Plus_Presentation_v2.pptx")
P.save(fp); os.remove("__ppt_temp.pptx")
print(f"DONE! Saved: {fp}")
print(f"Total slides: {len(P.slides)}")
